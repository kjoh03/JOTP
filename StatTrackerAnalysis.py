import numpy as np
import pandas as pd
import os
import pwd
import glob
import matplotlib
import matplotlib.pyplot as plt
from matplotlib.pyplot import figure
import datetime
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from datetime import date
import matplotlib.patches as mpatches
import StatTracker as ST

'''
def player_list(team):
    if team == 'Joy AC':
        playerlist = ['Joy AC Opponents','#2 Vukota Mastilovic','#3 Marco Corona Duran','#4, Noah Kantorowicz','#6 Martin Browne Jr','#8 Philip Caputo','#10 Whitney Browne','#11 Bennett Kouame','#12 Marshall Urban','#13 Luca Guarin','#14 Brian Kallman','#18 Dimitri Nair','#19 Mika Folstad','#22 David Riera','#24 Zinedine Kroeten','#28 Simeon Friesen','#29 Gabriel Eduarte','Henry Elias','Diego Paulin','Griffin Price','Jonathan Robles','#0 Tucker Mann','#1 Gage Steiner']
#2021        playerlist = ['Joy AC Opponents','#2 Devan DiGrado','#3 Marco Corona Duran','#4 Noah Kantorowicz','#5 Vukota Mastilovic','#6 David Riera','#7 Zinedine Kroeten','#8 Aiden Cavanaugh','#9 Darley Florvil','#10 Whitney Browne','#11 Emmanuel Iwe','#12 Marshall Urban','#14 Philip Caputo','#15 Otis Anderson','#16 Xavier Zengue','#17 Denilson Ramos','#18 Dennis Mensah','#19 Mika Folstad','#20 Jorge Radilla','#21 Luis Martinez Rojas','#22 Dimitri Nair','#23 Liam Vance','#24 Martin Browne Jr','#27 Abduselam Regassa','#28 Gabriel Eduarte','#29 Diego Paulin','#1 Dawson Fairchild','#0 Tucker Mann','#30 Ayuub Ahmed','Andrei Gotsmanov','Siddiq Madson-Keita','Henry Elias','Own Goal (Opponents)','Own Goal']
    elif team == 'Joy LaMancha':
        playerlist = ['Joy AC Opponents','Phil','Gabe','Henry','Mika','Luca','Siddiq','Bennett','Dimi','Andrew','Diego','Jack','Diego GK','Own Goal (Opponents)','Own Goal']
    elif team == 'Joy U19':
        playerlist = ['Joy AC Opponents','Adri','Henry','Michael','Phil','Mika','Victor','Elsini','Hassan','Oliver','Josiah','Noah','Sean','Isaiah','Carlitos','Yonas','Minoli','Gabe','Bennett','Sebe','Zekiah','Own Goal (Opponents)','Own Goal']
    elif team == 'Joy Kaghani':
        playerlist = ['Joy AC Opponents','Wilton','Si','Luca','Jack','Eric','Johnny','Manny','Max','Liam','Own Goal (Opponents)','Own Goal']
    elif team == 'Joy Tennessee Fainting':
        playerlist = ['Joy AC Opponents','Alanna','Dare','Allison','Iman','Riana','Sami','Aliviah','Sam','Ari','Caitlyn','Own Goal (Opponents)','Own Goal']
    elif team == 'Joy 06':
        playerlist = ['Joy AC Opponents','Nico', 'Tyler', 'Oliver', 'Isaiah', 'Will', 'Luca','Benji', 'Mikey', 'Kai', 'Eric', 'Makai', 'John', 'Leo', 'Lucas', 'Johnny','Jeremy','Own Goal (Opponents)','Own Goal'] 
    elif team == 'Cornell Mens Soccer':
        playerlist = ['Joy AC Opponents','#2 Gaurab Khadka','#3 Greg Pappadakis','#4 Emerson Roy','#5 Thomas Hamborg','#6 Drew Bruck','#7 Cian McNamara','#8 Nolan Zeger','#9 Aria Dehshid','#10 Owen Smith','#11 Andrew Lopez','#12 Galen Westervelt','#14 Sam Brueck','#15 Zach Miller','#16 Jonas Ricke','#17 Vance Wicker','#18 Mardoche Ntonku','#19 Justin Howe','#20 Brian Gin','#21 George Archer','#22 Blake Soto','#23 Bryce Scott','#24 Aron Mawia','#27 Federico Polidori','#28 Eddie Garces','#0 Will Bickel','#1 Jeremy Spina','#32 Mateo Ramirez']
    return playerlist
'''

def read_data(team,year):
    dfDict = {}
    username = pwd.getpwuid(os.getuid())[0]
    directory = '/Users/'+username+'/Desktop/Stat Tracker/'+team+'/'+year
    os.chdir(directory)
    username = pwd.getpwuid(os.getuid())[0]
    lis = glob.glob('*.csv')
    for file in lis:
        df = pd.read_csv(directory+'/'+file)
        dfDict[file.split('.')[0]] = df
    return dfDict

def read_all_data_to_one_df(dfDict,matchlist):
    singledf = pd.DataFrame()
    for match in matchlist:
        singledf = singledf.append(dfDict[match],ignore_index = True)
    return singledf

def team_dict(team,year):
    TeamDict = {}
    playerlist = ['Joy AC Opponents']+ST.player_list(team,year)[0][1:]
    for player in playerlist:
        PlayerDict = {}
        PlayerDict['Shots'] = 0
        PlayerDict['xG'] = 0
        PlayerDict['xA'] = 0
        PlayerDict['Goals'] = 0
        PlayerDict['Assists'] = 0
        PlayerDict['Shot Distance'] = 0
        TeamDict[player] = PlayerDict
    return TeamDict

def data_to_team_dict(matchdf,TeamDict):
    for row in range(len(matchdf.index)):
        if matchdf.loc[row,'Name'] == 'Opponent Player':
            TeamDict['Joy AC Opponents']['xG'] += matchdf.loc[row,'xG']
            if matchdf.loc[row,'Assist'] == 'Unassisted':
                pass
            else:
                TeamDict['Joy AC Opponents']['xA'] += matchdf.loc[row,'xG']
            if matchdf.loc[row,'Event'] == 'Goal' or matchdf.loc[row,'Event'] == 'Penalty Goal':
                TeamDict['Joy AC Opponents']['Goals'] += 1
                if matchdf.loc[row,'Assist'] == 'Unassisted':
                    pass
                else:
                    TeamDict['Joy AC Opponents']['Assists'] += 1
        else:
            TeamDict[matchdf.loc[row,'Name']]['xG'] += matchdf.loc[row,'xG']
            TeamDict[matchdf.loc[row,'Name']]['Shots'] += 1
            TeamDict[matchdf.loc[row,'Name']]['Shot Distance'] += matchdf.loc[row,'Shot Distance']
            if matchdf.loc[row,'Assist'] == 'Unassisted':
                pass
            else:
                TeamDict[matchdf.loc[row,'Assist']]['xA'] += matchdf.loc[row,'xG']                    
            if matchdf.loc[row,'Event'] == 'Goal' or matchdf.loc[row,'Event'] == 'Penalty Goal':
                TeamDict[matchdf.loc[row,'Name']]['Goals'] += 1
            if matchdf.loc[row,'Event'] == 'Goal' or matchdf.loc[row,'Event'] == 'Penalty Goal':
                if matchdf.loc[row,'Assist'] == 'Unassisted':
                    pass
                else:
                    TeamDict[matchdf.loc[row,'Assist']]['Assists'] += 1
    return TeamDict

def full_season_player_data(dfDict,TeamDict,matchlist):
    for match in matchlist:
        data_to_team_dict(dfDict[match],TeamDict)
    return TeamDict
    
def match_stats(dfDict,matchID):
    MatchDict = {}
    JoyDict = {}
    OpponentsDict = {}
    JoyDict['Shots'] = 0
    OpponentsDict['Shots'] = 0
    JoyDict['xG'] = 0
    OpponentsDict['xG'] = 0
    JoyDict['Goals'] = 0
    OpponentsDict['Goals'] = 0
    JoyDict['Score Progression'] = [[0,0,datetime.timedelta(minutes=0,seconds=0),'']]
    JoyDict['xG When Goal Scored'] = [0]
    JoyDict['Shot Chart Values'] = []
    OpponentsDict['Shot Chart Values'] = []
    JoyDict['Progressing xG Values'] = [0]
    OpponentsDict['Progressing xG Values'] = [0]
    JoyDict['Goals List'] = [0]
    OpponentsDict['Goals List'] = [0]
    JoyDict['Goals Time Stamps'] = [datetime.timedelta(minutes=0,seconds=0)]
    OpponentsDict['Goals Time Stamps'] = [datetime.timedelta(minutes=0,seconds=0)]
    JoyDict['Time Stamp List'] = [datetime.timedelta(minutes=0,seconds=0)]
    OpponentsDict['Time Stamp List'] = [datetime.timedelta(minutes=0,seconds=0)]
    JoyDict['Own Goals'] = []
    OpponentsDict['Own Goals'] = []
    match = dfDict[matchID]
    for row in range(len(match.index)):
        if match.loc[row,'Name'] == 'Opponent Player' and not match.loc[row,'Event'] == 'Own Goal' or not match.loc[row,'Name'] == 'Opponent Player' and match.loc[row,'Event'] == 'Own Goal':
            OpponentsDict['Shots'] += 1
            OpponentsDict['Shot Chart Values'].append((match.loc[row,'x'],match.loc[row,'y'],match.loc[row,'xG']))
            OpponentsDict['Progressing xG Values'].append(OpponentsDict['Progressing xG Values'][-1])
            OpponentsDict['Progressing xG Values'].append(OpponentsDict['Progressing xG Values'][-1]+match.loc[row,'xG'])
            OpponentsDict['Time Stamp List'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
            OpponentsDict['Time Stamp List'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
            OpponentsDict['xG'] += match.loc[row,'xG']
            if match.loc[row,'Event'] == 'Goal' or match.loc[row,'Event'] == 'Penalty Goal' or not match.loc[row,'Name'] == 'Opponent Player' and match.loc[row,'Event'] == 'Own Goal':
                OpponentsDict['Goals'] += 1
#                OpponentsDict['Goals List'].append(OpponentsDict['Goals List'][-1])
                OpponentsDict['Goals List'].append(OpponentsDict['Goals List'][-1]+1)
#                JoyDict['Goals List'].append(JoyDict['Goals List'][-1])
                OpponentsDict['Goals Time Stamps'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
#                OpponentsDict['Goals Time Stamps'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
                JoyDict['xG When Goal Scored'].append(OpponentsDict['Progressing xG Values'][-1])
                JoyDict['Score Progression'].append([JoyDict['Goals'],OpponentsDict['Goals'],OpponentsDict['Goals Time Stamps'][-1],matchID.split(' vs ')[1]])
        elif not match.loc[row,'Name'] == 'Opponent Player' and not match.loc[row,'Event'] == 'Own Goal':
            JoyDict['Shots'] += 1
            JoyDict['Shot Chart Values'].append((match.loc[row,'x'],match.loc[row,'y'],match.loc[row,'xG']))
            JoyDict['Progressing xG Values'].append(JoyDict['Progressing xG Values'][-1])
            JoyDict['Progressing xG Values'].append(JoyDict['Progressing xG Values'][-1]+match.loc[row,'xG'])
            JoyDict['Time Stamp List'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
            JoyDict['Time Stamp List'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
            JoyDict['xG'] += match.loc[row,'xG']
            if match.loc[row,'Event'] == 'Goal' or match.loc[row,'Event'] == 'Penalty Goal':
                JoyDict['Goals'] += 1
#                JoyDict['Goals List'].append(JoyDict['Goals List'][-1])
                JoyDict['Goals List'].append(JoyDict['Goals List'][-1]+1)
#                OpponentsDict['Goals List'].append(OpponentsDict['Goals List'][-1])
                JoyDict['Goals Time Stamps'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
#                JoyDict['Goals Time Stamps'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
                JoyDict['xG When Goal Scored'].append(JoyDict['Progressing xG Values'][-1])
                JoyDict['Score Progression'].append([JoyDict['Goals'],OpponentsDict['Goals'],JoyDict['Goals Time Stamps'][-1],'Joy AC'])
        elif match.loc[row,'Name'] == 'Opponent Player' and match.loc[row,'Event'] == 'Own Goal':
            JoyDict['Goals'] += 1
            JoyDict['Goals Time Stamps'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
            JoyDict['Score Progression'].append([JoyDict['Goals'],OpponentsDict['Goals'],JoyDict['Goals Time Stamps'][-1],'Joy AC'])
            JoyDict['xG When Goal Scored'].append(JoyDict['Progressing xG Values'][-1])
        elif not match.loc[row,'Name'] == 'Opponent Player' and match.loc[row,'Event'] == 'Own Goal':
            OpponentsDict['Shots'] += 1
            OpponentsDict['Goals Time Stamps'].append(datetime.timedelta(minutes=int(match.loc[row,'Time'].split(':')[0]), seconds=int(match.loc[row,'Time'].split(':')[1])))
            JoyDict['Score Progression'].append([JoyDict['Goals'],OpponentsDict['Goals'],OpponentsDict['Goals Time Stamps'][-1],matchID.split(' vs ')[1]])
            OpponentsDict['xG When Goal Scored'].append(OpponentsDict['Progressing xG Values'][-1])
    MatchDict['Joy AC'] = JoyDict
    MatchDict['Joy AC Opponents'] = OpponentsDict
    return MatchDict

def translate_times(times):
    newtimes = []
    for ele in times:
        totalsecs = ele.total_seconds()
        newtimes.append(totalsecs)
    return newtimes
    
def draw_xG_graphic(MatchDict,dfDict,matchID):
    match = dfDict[matchID]
    ### FIRST TEAM
    times = MatchDict['Joy AC']['Time Stamp List']
    newtimes = translate_times(times)
    xG = MatchDict['Joy AC']['Progressing xG Values']
    plt.plot(newtimes, xG,'r')
    ### SECOND TEAM
    times = MatchDict['Joy AC Opponents']['Time Stamp List']
    newtimes = translate_times(times)
    xG = MatchDict['Joy AC Opponents']['Progressing xG Values']
    plt.plot(newtimes, xG,'b')
    plt.legend([matchID.split(' vs ')[0],matchID.split(' vs ')[1]])

    ### GOALS FIRST TEAM
    times = []
    for ele in MatchDict['Joy AC']['Score Progression']:
        times.append(ele[2])
    newtimes = translate_times(times)
    for i in range(1,len(MatchDict['Joy AC']['Score Progression'])):
        if MatchDict['Joy AC']['Score Progression'][i][3] == 'Joy AC':
            plt.text(newtimes[i]-75, MatchDict['Joy AC']['xG When Goal Scored'][i]+.04,str(MatchDict['Joy AC']['Score Progression'][i][0])+'-'+str(MatchDict['Joy AC']['Score Progression'][i][1]), fontsize = 6, color = 'r')
        else:
            plt.text(newtimes[i]-75, MatchDict['Joy AC']['xG When Goal Scored'][i]+.04,str(MatchDict['Joy AC']['Score Progression'][i][0])+'-'+str(MatchDict['Joy AC']['Score Progression'][i][1]), fontsize = 6, color = 'b')
    '''
    ### GOALS FIRST TEAM
    times = MatchDict['Joy AC']['Goals Time Stamps']
    newtimes = translate_times(times)
    xG = MatchDict['Joy AC']['Goals List']
    plt.plot(newtimes, xG,'r',linestyle='dashed')
    ### GOALS SECOND TEAM
    times = MatchDict['Joy AC Opponents']['Goals Time Stamps']
    newtimes = translate_times(times)
    xG = MatchDict['Joy AC Opponents']['Goals List']
    plt.plot(newtimes, xG,'b',linestyle='dashed')
    '''

    timevalues = [0,600,1200,1800,2400,3000,3600,4200,4800,5400]
    values = ['0:00','10:00','20:00','30:00','40:00','50:00','60:00','70:00','80:00','90:00']
    plt.xticks(timevalues,values)
    plt.xticks(rotation = 45)
    plt.title("xG Progression")
    plt.xlabel("Game Time")
    plt.ylabel("Cumulative xG")
#    plt.savefig('/Users/keanjohansen/Desktop/JOTP/Joy AC Matches/xG Graphics/'+matchID, dpi = 300)
    plt.show()
    return ''
'''
def shot_chart(MatchDict,matchID):
    fig,ax = ST.draw_soccer_pitch(figsize=(9, 6.75))
    for shot in MatchDict['Joy AC']['Shot Chart Values']:
        circle = plt.Circle((shot[1],shot[0]),.25,color='red')   ### .25 RADIUS FOR OUTDOOR
        plt.text(shot[1]-1,shot[0]+1.5,str(round(shot[2],2)),fontsize = 6)  ### -1,+1.5
        ax.add_patch(circle)
    plt.title("Joy Shot Chart - "+str(MatchDict['Joy AC']['Shots'])+' Shots, '+str(round(MatchDict['Joy AC']['xG'],2))+' xG, '+str(MatchDict['Joy AC']['Goals'])+' Goals')
#    plt.savefig('/Users/keanjohansen/Desktop/JOTP/Joy AC Matches/Shot Charts/Joy AC Shots '+matchID, dpi = 300)
    fig2,ax2 = ST.draw_soccer_pitch(figsize=(9, 6.75))
    for shot in MatchDict['Joy AC Opponents']['Shot Chart Values']:
        circle = plt.Circle((shot[1],shot[0]),.25,color='red')   ### .25 RADIUS FOR OUTDOOR
        plt.text(shot[1]-1,shot[0]+1.5,str(round(shot[2],2)),fontsize = 6)  ### -1,+1.5
        ax2.add_patch(circle)
    plt.title("Opponents Shot Chart - "+str(MatchDict['Joy AC Opponents']['Shots'])+' Shots, '+str(round(MatchDict['Joy AC Opponents']['xG'],2))+' xG, '+str(MatchDict['Joy AC Opponents']['Goals'])+' Goals')
#    plt.savefig('/Users/keanjohansen/Desktop/JOTP/Joy AC Matches/Shot Charts/Joy AC Opponents Shots '+matchID, dpi = 300)
    plt.show()    
'''
def individual_performers_df(TeamDict):
    cols = ['Name','Shots','Goals','Assists','xG','xA','xG+xA','Avg Shot Dist']
    rowslist = []
    for player in TeamDict:
        if player == 'Joy AC Opponents':
            pass
        else:
            if TeamDict[player]['Shots'] == 0:
                rowslist.append([player,TeamDict[player]['Shots'],TeamDict[player]['Goals'],TeamDict[player]['Assists'],TeamDict[player]['xG'],TeamDict[player]['xA'],TeamDict[player]['xG']+TeamDict[player]['xA'],'NA'])
            else:
                rowslist.append([player,TeamDict[player]['Shots'],TeamDict[player]['Goals'],TeamDict[player]['Assists'],TeamDict[player]['xG'],TeamDict[player]['xA'],TeamDict[player]['xG']+TeamDict[player]['xA'],round(TeamDict[player]['Shot Distance']/TeamDict[player]['Shots'],2)])
    df = pd.DataFrame(rowslist,columns = cols)
    df = df.sort_values(by=['xG+xA','xG','xA'],ascending = False).reset_index(drop=True)
    print(df)
    df.to_csv('/Users/keanjohansen/Desktop/JOTP/Outdated/test.csv',index=False,header=True)
    return df

def shot_chart(shotdf,save):
    shotcounter = 0
    xGcounter = 0
    goalcounter = 0
    fig,ax = ST.draw_soccer_pitch(figsize=(9, 6.75)) ### CHANGE INDOOR/OUTDOOR
    ### SHOT CHART
    for row in range(len(shotdf)):
        if not shotdf.loc[row,'Name'] == 'Opponent Player' and shotdf.loc[row,'Event'] == 'Own Goal':
            goalcounter +=1
            colour = 'lime'
            circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color=colour)  ### .25 RADIUS FOR OUTDOOR
            plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6)     ### -1,+1.5
            ax.add_patch(circle)
        if shotdf.loc[row,'Name'] == 'Opponent Player' and not shotdf.loc[row,'Event'] == 'Own Goal':
            if shotdf.loc[row,'Event'] == 'Goal' or shotdf.loc[row,'Event'] == 'Penalty Goal' or shotdf.loc[row,'Event'] == 'Own Goal':
                goalcounter +=1
                colour = 'lime'
            else:
                colour = 'red'
            shotcounter += 1
            xGcounter += shotdf.loc[row,'xG']
            circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color=colour)  ### .25 RADIUS FOR OUTDOOR
            plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6)     ### -1,+1.5
            ax.add_patch(circle)
    green_patch = mpatches.Patch(color='lime', label='Goal')
    red_patch = mpatches.Patch(color='red', label='Shot')
    blue_patch = mpatches.Patch(color='blue', label='Shot Assist')
    plt.legend(handles=[green_patch,red_patch,blue_patch],loc = 'lower left')            
    plt.title("Opponents Shot Chart - "+str(shotcounter)+' Shots, '+str(round(xGcounter,2))+' xG, '+str(goalcounter)+' Goals')
    if save == 'yes':
        plt.savefig('/Users/keanjohansen/Desktop/JOTP/Graphics/Joy Opponents', dpi = 300)
    fig2,ax2 = ST.draw_soccer_pitch(figsize=(9, 6.75))
    shotcounter = 0
    xGcounter = 0
    goalcounter = 0
    for row in range(len(shotdf)):
        if shotdf.loc[row,'Name'] == 'Opponent Player' and shotdf.loc[row,'Event'] == 'Own Goal':
            goalcounter +=1
            colour = 'lime'
            circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color=colour)  ### .25 RADIUS FOR OUTDOOR
            plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6)     ### -1,+1.5
            ax2.add_patch(circle)
        if not shotdf.loc[row,'Name'] == 'Opponent Player' and not shotdf.loc[row,'Event'] == 'Own Goal':
            if shotdf.loc[row,'Event'] == 'Goal' or shotdf.loc[row,'Event'] == 'Penalty Goal' or shotdf.loc[row,'Event'] == 'Own Goal':
                goalcounter +=1
                colour = 'lime'
            else:
                colour = 'red'
            shotcounter += 1
            xGcounter += shotdf.loc[row,'xG']
            circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color=colour)  ### .25 RADIUS FOR OUTDOOR
            plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6)     ### -1,+1.5
            ax2.add_patch(circle)
    plt.title("Cornell Shot Chart - "+str(shotcounter)+' Shots, '+str(round(xGcounter,2))+' xG, '+str(goalcounter)+' Goals')
    if save == 'yes':
        plt.savefig('/Users/keanjohansen/Desktop/JOTP/Graphics/Joy', dpi = 300)
    green_patch = mpatches.Patch(color='lime', label='Goal')
    red_patch = mpatches.Patch(color='red', label='Shot')
    blue_patch = mpatches.Patch(color='blue', label='Shot Assist')
    plt.legend(handles=[green_patch,red_patch,blue_patch],loc = 'lower left')
    if not save == 'yes':
        plt.show()

def individual_player_shot_charts(player,shotdf,save):
    shotcounter = 0
    xGcounter = 0
    goalcounter = 0
    fig,ax = ST.draw_soccer_pitch(figsize=(9, 6.75)) ### CHANGE INDOOR/OUTDOOR
    ### SHOT CHART
    for row in range(len(shotdf)):
        if shotdf.loc[row,'Name'] == player:
            if shotdf.loc[row,'Event'] == 'Goal' or shotdf.loc[row,'Event'] == 'Penalty Goal':
                goalcounter +=1
                colour = 'lime'
            else:
                colour = 'red'
            shotcounter += 1
            xGcounter += shotdf.loc[row,'xG']
            circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color=colour)  ### .25 RADIUS FOR OUTDOOR
            plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6)     ### -1,+1.5
            ax.add_patch(circle)
    ### ASSIST CHART
    shotassistcounter = 0
    xAcounter = 0
    assistcounter = 0
    for row in range(len(shotdf)):
        if shotdf.loc[row,'Assist'] == player:
            if shotdf.loc[row,'Event'] == 'Goal' or shotdf.loc[row,'Event'] == 'Penalty Goal':
                assistcounter +=1
            shotassistcounter += 1
            xAcounter += shotdf.loc[row,'xG']
            circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color='blue')  ### .25 RADIUS FOR OUTDOOR
            plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6) ### -1,+1.5
            ax.add_patch(circle)
    plt.title(player+' Shot Chart'+'\n'+str(shotcounter)+' Shots, '+str(round(xGcounter,2))+' xG, '+str(goalcounter)+' Goals, '+'\n'+str(shotassistcounter)+' Shot Assists, '+str(round(xAcounter,2))+' xA, '+str(assistcounter)+' Assists')
    green_patch = mpatches.Patch(color='lime', label='Goal')
    red_patch = mpatches.Patch(color='red', label='Shot')
    blue_patch = mpatches.Patch(color='blue', label='Shot Assist')
    plt.legend(handles=[green_patch,red_patch,blue_patch],loc = 'lower left')
    if save == 'yes':
        plt.savefig('/Users/keanjohansen/Desktop/JOTP/Graphics/'+player, dpi = 300)
    else:
        plt.show()    

def xG_xA_player_graphic(PlayerDict,save):
    for player in PlayerDict:
        if player == 'Joy AC Opponents':
            pass
        elif PlayerDict[player]['xG'] == 0 and PlayerDict[player]['xA'] == 0:
            pass
        else:
            plt.scatter(PlayerDict[player]['xG'],PlayerDict[player]['xA'],label = player)
            plt.annotate(player,(PlayerDict[player]['xG'],PlayerDict[player]['xA']+.04),ha='center',fontsize = 6)
            plt.title("xG xA Graphic")
            plt.xlabel("Expected Goals")
            plt.ylabel("Expected Assists")
#    plt.plot((1,0),(0,1),'r--')
#    plt.plot((2,0),(0,2),'r--')
    if save == 'yes':
        plt.savefig('/Users/keanjohansen/Desktop/JOTP/Graphics/xGxAGraphic', dpi = 300)
    else:
        plt.show()

def shot_distance(x,y):
    distance = math.sqrt(x**2+y**2)
    return distance


def season_progression_graphic(player,dfDict,matchlist):
    xGvalues = []
    xAvalues = []
    xG = 0      #should be able to move these (xG and xA) one loop down to switch graph type
    xA = 0
    xaxis = []
    i = 1
    for match in matchlist:
        for row in range(len(dfDict[match])):
            if dfDict[match].loc[row,'Name'] == player:
                xG += dfDict[match].loc[row,'xG']
            elif dfDict[match].loc[row,'Assist'] == player:
                xA += dfDict[match].loc[row,'xG']
        xGvalues.append(xG)
        xAvalues.append(xA)
        xaxis.append(i)
        i+=1
    plt.plot(xaxis,xGvalues)
    plt.plot(xaxis,xAvalues)
    plt.legend(['xG','xA'])
    plt.title(player+'`s Season Progression of xG, xA')
    plt.xlabel('Match Number')
    plt.ylabel('Expected Value')
    plt.xticks(xaxis,['Luther','Simpson','Rockford','N. Wesleyan','Bethel','Coe','Knox','Central','Ripon','Grinnell','Buena Vista','IC','LF'])
    plt.xticks(rotation = 45,fontsize = 7)
    plt.show()



def create_presentation(PlayerDict,filename,team,matchlist,year):
    xG_xA_player_graphic(PlayerDict,'yes')
    prs = Presentation()
    ### TITLE SLIDE
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = team+' Season Stats'
    subtitle.text = str(date.today().strftime("%d/%m/%Y"))
    ### SLIDE ONE
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    title1 = slide.shapes.title
    title1.text = team + ' Season Shot Chart'
    dfDict = read_data(team,year)
    shotdf = read_all_data_to_one_df(dfDict,matchlist)
    shot_chart(shotdf,'yes')
    left = Inches(.75)
    top = Inches(1.5)
    height = width = Inches(6.5)
    img_path = '/Users/keanjohansen/Desktop/JOTP/Graphics/Joy.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    ### SLIDE TWO
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    title1 = slide.shapes.title
    title1.text = team + ' Opponents Season Shot Chart'
    left = Inches(.75)
    top = Inches(1.5)
    height = width = Inches(6.5)
    img_path = '/Users/keanjohansen/Desktop/JOTP/Graphics/Joy Opponents.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    ### SLIDE THREE
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    title1 = slide.shapes.title
    title1.text = 'Player Statistics'
    left = Inches(.75)
    top = Inches(1)
    height = width = Inches(6.5)
    img_path = '/Users/keanjohansen/Desktop/JOTP/Graphics/xGxAGraphic.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    ### PLAYER SLIDES
    TeamDict = team_dict(team,year)
    PlayerDict = full_season_player_data(dfDict,TeamDict,matchlist)
    for player in sorted(PlayerDict.items(), key=lambda x: x[1]['xG'], reverse=True):
        if player[0] == 'Joy AC Opponents':
            print(player[0])
            pass
        else:
            print(player)
            individual_player_shot_charts(player[0],shotdf,'yes')
            plt.close()
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)
            title1 = slide.shapes.title
            title1.text = player[0] + ' Season Stats'
            left = Inches(-1)
            top = Inches(1.5)
            height = width = Inches(6.5)
            img_path = '/Users/keanjohansen/Desktop/JOTP/Graphics/'+player[0]+'.png'
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
            text2 = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(2), Inches(2))
            text2.text_frame.text = ''       
    prs.save(filename)




def unassisted_shot_chart(team,shotdf):
    shotcounter = 0
    xGcounter = 0
    goalcounter = 0
    fig,ax = ST.draw_soccer_pitch(figsize=(9, 6.75)) ### CHANGE INDOOR/OUTDOOR
    ### SHOT CHART
    if team == 'Opponents':
        for row in range(len(shotdf)):
            if shotdf.loc[row,'Assist'] == 'Unassisted':
                if shotdf.loc[row,'Name'] == 'Opponent Player':
                    if shotdf.loc[row,'Event'] == 'Goal' or shotdf.loc[row,'Event'] == 'Penalty Goal' or shotdf.loc[row,'Event'] == 'Own Goal' and not shotdf.loc[row,'Name'] == 'Opponent Player':
                        goalcounter +=1
                        colour = 'lime'
                    else:
                        colour = 'red'
                    shotcounter += 1
                    xGcounter += shotdf.loc[row,'xG']
                    circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color=colour)  ### .25 RADIUS FOR OUTDOOR
                    plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6)     ### -1,+1.5
                    ax.add_patch(circle)
    else:
        for row in range(len(shotdf)):
            if shotdf.loc[row,'Assist'] == 'Unassisted':
                if not shotdf.loc[row,'Name'] == 'Opponent Player':
                    if shotdf.loc[row,'Event'] == 'Goal' or shotdf.loc[row,'Event'] == 'Penalty Goal' or shotdf.loc[row,'Event'] == 'Own Goal' and shotdf.loc[row,'Name'] == 'Opponent Player':
                        goalcounter +=1
                        colour = 'lime'
                    else:
                        colour = 'red'
                    shotcounter += 1
                    xGcounter += shotdf.loc[row,'xG']
                    circle = plt.Circle((shotdf.loc[row,'y'],shotdf.loc[row,'x']),.25,color=colour)  ### .25 RADIUS FOR OUTDOOR
                    plt.text(shotdf.loc[row,'y']-1,shotdf.loc[row,'x']+1.5,str(round(shotdf.loc[row,'xG'],2)),fontsize = 6)     ### -1,+1.5
                    ax.add_patch(circle)
    plt.title(team+'`s Unassisted Shot Chart'+'\n'+str(shotcounter)+' Shots, '+str(round(xGcounter,2))+' xG, '+str(goalcounter)+' Goals, '+'\n'+str(shotcounter)+' Shot Assists')
    green_patch = mpatches.Patch(color='lime', label='Goal')
    red_patch = mpatches.Patch(color='red', label='Shot')
    plt.legend(handles=[green_patch,red_patch],loc = 'lower left')
    plt.show() 

    
if __name__ == '__main__':
    year = '2022'
    schedule = ['Cornell Mens Soccer vs Luther','Cornell Mens Soccer vs Simpson','Cornell Mens Soccer vs Rockford','Cornell Mens Soccer vs Nebraska Wesleyan','Cornell Mens Soccer vs Bethel','Cornell Mens Soccer vs Coe','Cornell Mens Soccer vs Knox','Cornell Mens Soccer vs Central','Cornell Mens Soccer vs Ripon','Cornell Mens Soccer vs Grinnell','Cornell Mens Soccer vs Buena Vista','Cornell Mens Soccer vs Illinois College','Cornell Mens Soccer vs Lake Forest']
    #schedule = ['Cornell vs Luther','Cornell vs Coe','Cornell vs Rockford','Cornell vs Fontbonne','Cornell vs Simpson','Cornell vs Nebraska Wesleyan','Cornell vs Central','Cornell vs Monmouth','Cornell vs Coe (away)','Cornell vs Ripon','Cornell vs Iowa Wesleyan','Cornell vs Grinnell','Cornell vs Illinois College','Cornell vs Lake Forest','Cornell vs Knox','Cornell vs Lawrence']
    
    dfDict = read_data('Cornell Mens Soccer',year)
    TeamDict = team_dict('Cornell Mens Soccer',year)
    PlayerDict = full_season_player_data(dfDict,TeamDict,schedule)

    shotdf = read_all_data_to_one_df(dfDict,schedule)
    unassisted_shot_chart('Opponents',shotdf)

#    create_presentation(PlayerDict,'/Users/keanjohansen/Desktop/JOTP/Outdated/Cornell Mens Soccer test.pptx','Cornell Mens Soccer', schedule,year)
    '''
    dfDict = read_data('Joy AC')   
    TeamDict = team_dict('Joy AC')
    PlayerDict = full_season_player_data(dfDict,TeamDict,['Match Day 1 Joy vs Fusion','Match Day 2 Joy vs La Crosse','Match Day 4 Joy vs Med City','Match Day 5 Joy vs Sioux Falls','Match Day 7 Joy vs MPLS City','Match Day 12 Joy vs MPLS City'])

    create_presentation(PlayerDict,'/Users/keanjohansen/Desktop/JOTP/Outdated/test.pptx','Joy AC',matchlist = ['Match Day 1 Joy vs Fusion','Match Day 2 Joy vs La Crosse','Match Day 4 Joy vs Med City','Match Day 5 Joy vs Sioux Falls','Match Day 7 Joy vs MPLS City','Match Day 12 Joy vs MPLS City'])
    '''
    
#    season_progression_graphic('#13 Zach Miller',dfDict,schedule)

#    xG_xA_player_graphic(PlayerDict)
#    print(PlayerDict)
#    print(PlayerDict.keys())
 #   print(sorted(PlayerDict.items(), key=lambda x: x[1]['xG'], reverse=True))

    
#    shotdf = read_all_data_to_one_df(dfDict,['Match Day 1 Joy vs Fusion','Match Day 2 Joy vs La Crosse','Match Day 4 Joy vs Med City','Match Day 5 Joy vs Sioux Falls','Match Day 7 Joy vs MPLS City','Match Day 12 Joy vs MPLS City'])

#    shot_chart(shotdf,'no')

    
#    print(full_season_player_data(dfDict,TeamDict,['Match Day 1 Joy vs Fusion','Match Day 2 Joy vs La Crosse','Match Day 4 Joy vs Med City','Match Day 5 Joy vs Sioux Falls','Match Day 7 Joy vs MPLS City','Match Day 12 Joy vs MPLS City']))
#    MatchDict = match_stats(dfDict,'U19s Joy vs Princeton Youth')
#    draw_xG_graphic(MatchDict,dfDict,'Match Day 5 Joy vs Sioux Falls')

#    individualmatchplayerstats = data_to_team_dict(dfDict['U19s Joy vs Princeton Youth'],TeamDict)
#    individual_performers_df(TeamDict)
#    fullseason = full_season_player_data(dfDict,TeamDict,['U19s Joy vs Princeton Youth'])
#    individual_performers_df(TeamDict)

#    shotdf = read_all_data_to_one_df(dfDict,['Match Day 1 Joy vs Fusion','Match Day 5 Joy vs Sioux Falls'])
#    individual_player_shot_charts('#11 Emmanuel Iwe',shotdf,'no')
    

#    playerlist = ['Joy AC Opponents','#2 Devan DiGrado','#3 Marco Corona Duran','#4 Noah Kantorowicz','#5 Vukota Mastilovic','#6 David Riera','#7 Zinedine Kroeten','#8 Aiden Cavanaugh','#9 Darley Florvil','#10 Whitney Browne','#11 Emmanuel Iwe','#12 Marshall Urban','#14 Philip Caputo','#15 Otis Anderson','#16 Xavier Zengue','#17 Denilson Ramos','#18 Dennis Mensah','#19 Mika Folstad','#20 Jorge Radilla','#21 Luis Martinez Rojas','#22 Dimitri Nair','#23 Liam Vance','#24 Martin Browne Jr','#27 Abduselam Regassa','#28 Gabriel Eduarte','#29 Diego Paulin','#1 Dawson Fairchild','#0 Tucker Mann','#30 Ayuub Ahmed','Andrei Gotsmanov','Siddiq Madson-Keita','Henry Elias']
#    playerlist = ['Joy AC Opponents','Phil','Gabe','Henry','Mika','Luca','Siddiq','Bennett','Dimi','Andrew','Diego','Jack','Diego GK']
#    playerlist = ['Joy AC Opponents','Wilton','Si','Luca','Jack','Eric','Johnny','Manny','Max','Liam']
#    playerlist = ['Joy AC Opponents','Alanna','Dare','Allison','Iman','Riana','Sami','Aliviah','Sam','Ari','Caitlyn']
#    playerlist = ['Joy AC Opponents','Adri','Henry','Michael','Victor','Elsini','Hassan','Oliver','Josiah','Noah','Sean','Isaiah','Carlitos','Yonas','Minoli','Gabe','Bennett','Sebe','Zekiah']
#    playerlist = ['Joy AC Opponents','Own Goal','Own Goal (Opponents)','Nico', 'Tyler', 'Oliver', 'Isaiah', 'Will', 'Luca','Benji', 'Mikey', 'Kai', 'Eric', 'Makai', 'John', 'Leo', 'Lucas', 'Johnny','Jeremy'] 


    
